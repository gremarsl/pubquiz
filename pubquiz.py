import os

import six

from classes.FileReader import FileReader
from pptx import Presentation
#to copy slides from other presentation
import copy
from classes.layout import Layout


def duplicate_slide(pres, index):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[1]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(value.reltype,
                                                    value._target,
                                                    value.rId)

    return copied_slide

def copy_slide_from_external_prs(prs):

    # copy from external presentation all objects into the existing presentation
    #external_pres = Presentation("/home/marcel/PycharmProjects/pubquiz/classes/layout.pptx")
    external_pres = Presentation("/home/marcel/PycharmProjects/pubquiz/Template.pptx")

    # specify the slide you want to copy the contents from
    ext_slide = external_pres.slides[1]

    # Define the layout you want to use from your generated pptx
    SLD_LAYOUT = 0
    slide_layout = prs.slide_layouts[SLD_LAYOUT]

    # create now slide, to copy contents to
    curr_slide = prs.slides.add_slide(slide_layout)

    # now copy contents from external slide, but do not copy slide properties
    # e.g. slide layouts, etc., because these would produce errors, as diplicate
    # entries might be generated
    for shp in ext_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return prs

def main():
    filereader = FileReader()
    directoryToFile = os.path.dirname(__file__)
    print(directoryToFile)
    #list of questions with question objects and its attributes
    questions = filereader.readAllCsvFiles(directoryToFile)
    #layout = Layout()
    #titleslide=layout.getTitleSlide()
    #SLD_LAYOUT_TITLE_AND_CONTENT = 0

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    #print(slide_layout)
    slide = prs.slides.add_slide(slide_layout)
    #copy_slide_from_external_prs(prs)
    copy_slide_from_external_prs(prs)

    '''
    titleSlide = layout.getTitleSlide()
    slide_layout = presentation.slide_layouts[layout.getTitleSlide()]
    presentation.slides.add_slide(titleSlide)
    '''
    prs.save('pubquiz.pptx')

if __name__ == '__main__':
    main()

