import copy
import six
from pptx import Presentation
from enum import IntEnum

# layout file name
fileName_layout = "classes/layout.pptx"

# ENUM used to define where are the layouts in the pptx
class SlideIndex(IntEnum):
    TITLE_SLIDE = 0
    SOG_AD = 1
    RULES = 2
    ROUND_TITLE = 3
    QUESTION_OPEN = 4
    QUESTION_ALTERNATIVE = 5
    QUESTION_MIDDLE = 6
    QUESTION_MIDDLE_ANSWER = 7
    RESULTS = 7



# the Layout Class
class Layout:
    """
    Used to get the Slides objects used as template
    """

    def __init__(self, pathLayout=fileName_layout):
        # opening the layout
        self.prs = Presentation(fileName_layout)

    def getTitleSlide(self):
        """
        Get the slide used as the Title one, the first
        :return: Slide object
        """
        return self.prs.slides[SlideIndex.TITLE_SLIDE]

    def getSOGadvertisement(self):
        """
        Get the slide used to show what is SOG
        :return: Slide object
        """
        return self.prs.slides[SlideIndex.SOG_AD]

    def getGameRules(self):
        """
        Get the slide with the game rules
        :return: Slide Object
        """
        return self.prs.slides[SlideIndex.RULES]
    def getQuestionOpen(self):
        """
        Get the slide with a layout for open (not multiple choice) question
        :return: Slide Object
        """
        return self.prs.slides[SlideIndex.QUESTION_OPEN]
    def getQuestionMultipleChoices(self):
        """
        Get the slide with a layout for multiple choice questions
        :return: Slide Object
        """
        return self.prs.slides[SlideIndex.QUESTION_ALTERNATIVE]
    def getRoundTitle(self):
        """
        The slide used at the beginning of every round
        :return: Slide Object
        """
        return self.prs.slides[SlideIndex.ROUND_TITLE]
    def getQuestionMiddle(self):
        """
        The slide used in-between the rounds (Zwischenfrage)
        :return: Slide Object
        """
        return self.prs.slides[SlideIndex.QUESTION_MIDDLE]
    def getQuestionMiddleAnswer(self):
        """
        The slide used to show the answer for the in-between questions
        :return: Slide Object
        """
        return self.prs.slides[SlideIndex.QUESTION_MIDDLE_ANSWER]
    def getResults(self):
        """
        The slide used to show the charts with results
        :return: Slide Object
        """
        return self.prs.slides[SlideIndex.RESULTS]

    def duplicate_slide(self, pres, index):
        template = pres.slides[index]
        try:
            blank_slide_layout = pres.slide_layouts[12]
        except:
            blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

        copied_slide = pres.slides.add_slide(blank_slide_layout)

        for shp in template.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for _, value in six.iteritems(template.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                copied_slide.part.rels.add_relationship(value.reltype,
                                                        value._target,
                                                        value.rId)

        return copied_slide