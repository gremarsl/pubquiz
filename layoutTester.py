#!/usr/bin/python3

# testing our Layout class
from classes.layout import Layout
from pptx import Presentation
from pptx.util import Inches, Pt

# for testing purposes
fileName_output = "layoutTester.pptx"

layout = Layout()

prs = Presentation("Template.pptx")

slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

left = Inches(2.5)
top = Inches(1)
width = Inches(8)
height = Inches(5.5)

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "SOG Pubquiz"

#prs.slides[0] = layout.getGameRules()

prs.save(fileName_output)
